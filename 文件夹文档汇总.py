import os
import tkinter as tk
from tkinter import filedialog, messagebox
from pathlib import Path
from typing import List, Dict, Any, Optional

try:
    import xlrd
except ImportError:
    xlrd = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


def install_dependencies():
    import subprocess
    import sys
    global xlrd, openpyxl, Document, Presentation
    packages = []
    if xlrd is None:
        packages.append('xlrd')
    if openpyxl is None:
        packages.append('openpyxl')
    if Document is None:
        packages.append('python-docx')
    if Presentation is None:
        packages.append('python-pptx')
    if packages:
        print(f"正在安装缺少的依赖: {packages}")
        subprocess.check_call([sys.executable, '-m', 'pip', 'install'] + packages)
        import xlrd
        import openpyxl
        from docx import Document
        from pptx import Presentation


def read_xls(file_path: str) -> List[List[str]]:
    workbook = xlrd.open_workbook(file_path)
    data = []
    for sheet in workbook.sheets():
        sheet_data = []
        for row_idx in range(sheet.nrows):
            row_data = []
            for col_idx in range(sheet.ncols):
                cell = sheet.cell(row_idx, col_idx)
                row_data.append(str(cell.value))
            sheet_data.append(row_data)
        data.append(sheet_data)
    return data


def read_xlsx(file_path: str) -> List[List[str]]:
    workbook = openpyxl.load_workbook(file_path, data_only=True)
    data = []
    for sheet in workbook.worksheets:
        sheet_data = []
        for row in sheet.iter_rows(values_only=True):
            row_data = [str(cell) if cell is not None else '' for cell in row]
            sheet_data.append(row_data)
        data.append(sheet_data)
    return data


def convert_xls_to_markdown_table(data: List[List[str]], file_name: str) -> str:
    md_content = []

    for sheet_idx, sheet_data in enumerate(data):
        if len(data) > 1:
            md_content.append(f"\n### Sheet {sheet_idx + 1}\n")

        if not sheet_data:
            continue

        col_count = len(sheet_data[0])
        md_content.append("| " + " | ".join([f"Column {i+1}" for i in range(col_count)]) + " |")
        md_content.append("| " + " | ".join(["---" for _ in range(col_count)]) + " |")

        for row in sheet_data:
            while len(row) < col_count:
                row.append('')
            md_content.append("| " + " | ".join(row) + " |")

    return "\n".join(md_content)


def read_docx(file_path: str) -> str:
    doc = Document(file_path)
    content = []
    for para in doc.paragraphs:
        content.append(para.text)
    return "\n".join(content)


def read_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    content = []
    for slide_idx, slide in enumerate(prs.slides):
        content.append(f"\n--- Slide {slide_idx + 1} ---\n")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content.append(shape.text)
    return "\n".join(content)


def get_project_tree(selected_dir: str) -> List[str]:
    structure = []
    for root, dirs, files in os.walk(selected_dir):
        rel_path = os.path.relpath(root, selected_dir)
        level = rel_path.count(os.sep) if rel_path != '.' else 0
        indent = '│   ' * (level) + '├── ' if level > 0 else ''
        structure.append(f"{indent}{os.path.basename(root)}/")

        sub_indent = '│   ' * (level + 1) + '├── '
        for f in files:
            structure.append(f"{sub_indent}{f}")
    return structure


def read_text_file(file_path: Path) -> Optional[str]:
    if file_path.stat().st_size > 10 * 1024 * 1024:
        return None
    try:
        with open(file_path, 'rb') as in_f:
            content = in_f.read()
            if b'\x00' in content:
                return None
            text_content = content.decode('utf-8')
        return text_content
    except (UnicodeDecodeError, PermissionError):
        return None
    except Exception:
        return None


def process_xls_file(file_path: Path) -> Optional[str]:
    ext = file_path.suffix.lower()
    try:
        if ext == '.xls':
            if xlrd is None:
                install_dependencies()
            data = read_xls(str(file_path))
        elif ext == '.xlsx':
            if openpyxl is None:
                install_dependencies()
            data = read_xlsx(str(file_path))
        else:
            return None

        return convert_xls_to_markdown_table(data, file_path.name)
    except Exception:
        return None


def process_docx_file(file_path: Path) -> Optional[str]:
    try:
        if Document is None:
            install_dependencies()
        return read_docx(str(file_path))
    except Exception:
        return None


def process_pptx_file(file_path: Path) -> Optional[str]:
    try:
        if Presentation is None:
            install_dependencies()
        return read_pptx(str(file_path))
    except Exception:
        return None


def generate_folder_documentation(selected_dir: str, output_file: str):
    folder_name = os.path.basename(selected_dir)
    md_content = []

    md_content.append(f"# {folder_name}\n")
    md_content.append(f"\n> 文件夹路径: {selected_dir}\n")

    md_content.append(f"\n## 项目结构\n")
    tree_lines = get_project_tree(selected_dir)
    md_content.append("```\n" + "\n".join(tree_lines) + "\n```")

    target_text_ext = {
        '.txt', '.md', '.json', '.xml', '.html', '.htm',
        '.c', '.cpp', '.h', '.hpp', '.java', '.py', '.js', '.php',
        '.rb', '.swift', '.go', '.kt', '.scala', '.cs', '.ts', '.tsx',
        '.vue', '.jsx', '.r', '.pl', '.lua', '.sh', '.bash',
        '.ini', '.yaml', '.yml', '.conf', '.cfg', '.properties',
        '.css', '.scss', '.less', '.sql', '.gitignore'
    }

    xls_ext = {'.xls', '.xlsx'}
    doc_ext = {'.docx'}
    ppt_ext = {'.pptx'}

    for root, _, files in os.walk(selected_dir):
        for file in files:
            file_path = Path(root) / file
            rel_path = file_path.relative_to(selected_dir)

            if file_path.suffix.lower() in xls_ext:
                md_content.append(f"\n\n## 文件: {rel_path}\n")
                xls_md = process_xls_file(file_path)
                if xls_md:
                    md_content.append(xls_md)
                else:
                    md_content.append(f"\n> 无法读取此文件\n")
            elif file_path.suffix.lower() in doc_ext:
                md_content.append(f"\n\n## 文件: {rel_path}\n")
                doc_content = process_docx_file(file_path)
                if doc_content:
                    md_content.append(f"\n```\n{doc_content}\n```")
                else:
                    md_content.append(f"\n> 无法读取此文件\n")
            elif file_path.suffix.lower() in ppt_ext:
                md_content.append(f"\n\n## 文件: {rel_path}\n")
                ppt_content = process_pptx_file(file_path)
                if ppt_content:
                    md_content.append(f"\n```\n{ppt_content}\n```")
                else:
                    md_content.append(f"\n> 无法读取此文件\n")
            elif file_path.suffix.lower() in target_text_ext:
                text_content = read_text_file(file_path)
                if text_content:
                    md_content.append(f"\n\n## 文件: {rel_path}\n")
                    md_content.append(f"\n```\n{text_content}\n```")

    with open(output_file, 'w', encoding='utf-8') as f:
        f.write("\n".join(md_content))


def select_folder() -> Optional[str]:
    selected_dir = filedialog.askdirectory(
        title='请选择要汇总的文件夹',
        initialdir=os.path.dirname(os.path.abspath(__file__))
    )
    return selected_dir if selected_dir else None


def main():
    root = tk.Tk()
    root.withdraw()

    selected_dir = select_folder()
    if not selected_dir:
        messagebox.showinfo("提示", "操作已取消")
        return

    try:
        folder_name = os.path.basename(selected_dir)
        output_file = os.path.join(os.path.dirname(os.path.abspath(__file__)), f"{folder_name}.md")

        generate_folder_documentation(selected_dir, output_file)

        messagebox.showinfo("成功", f"文档已生成:\n{output_file}")

    except Exception as e:
        messagebox.showerror("错误", f"发生错误：{str(e)}")
        import traceback
        traceback.print_exc()


if __name__ == "__main__":
    main()
